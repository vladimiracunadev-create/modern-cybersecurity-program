# -*- coding: utf-8 -*-
"""
Genera, POR PARTE, el material descargable de cada clase:
  - una guía en PDF (render del README vía Microsoft Edge headless)
  - una presentación PPTX (python-pptx) resumida por secciones, con una
    diapositiva por diagrama de la clase
y añade a cada README una sección "📥 Material descargable" con los enlaces.

Los bloques ```mermaid llegan ya dibujados desde scripts/mermaid_svg.py, que
los cachea por hash: el PDF incrusta el SVG y la presentación su versión en PNG,
porque python-pptx solo sabe insertar imágenes de mapa de bits. Sin eso el PDF mostraba el codigo del diagrama en crudo
justo donde deberia estar el grafico; y dejando que mermaid.js dibujara durante
el print-to-pdf, el resultado dependia de si la CDN respondia a tiempo. Ademas,
un diagrama con la sintaxis rota produce un SVG con el cartel de "Syntax error":
la cache lo rechaza, asi que no se cuela en el material como si fuera contenido.

Uso:  python scripts/generar_material.py <indice_de_parte>     # p. ej. 0, 1, 2 ...
      python scripts/generar_material.py <idx> --solo-una      # solo la 1ª clase (prueba)
      python scripts/generar_material.py --todas               # todas las partes seguidas

Requisitos: python-pptx, markdown, y Microsoft Edge (o Chrome) instalado.
"""
from __future__ import annotations
import glob
import html as htmllib
import os
import re
import subprocess
import sys
import tempfile

import markdown
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from mermaid_svg import dibujar as dibujar_diagramas  # noqa: E402
from mermaid_svg import rasterizar as rasterizar_diagramas  # noqa: E402
from salida_atomica import publicar  # noqa: E402

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CLASSES = os.path.join(ROOT, "classes")

NAVEGADORES = [
    r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
    r"C:\Program Files\Google\Chrome\Application\chrome.exe",
    r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
]


def navegador() -> str:
    for p in NAVEGADORES:
        if os.path.isfile(p):
            return p
    raise SystemExit("No se encontró Edge ni Chrome para generar PDFs.")


CSS = """
@page { size: A4; margin: 16mm 15mm; }
* { box-sizing: border-box; }
body { font-family: 'Segoe UI', system-ui, Arial, sans-serif; font-size: 10.5pt;
       line-height: 1.5; color: #14181d; }
h1 { font-size: 20pt; color: #0b3d2e; border-bottom: 2px solid #2e8b57; padding-bottom: 4px; }
h2 { font-size: 13.5pt; color: #0b3d2e; margin-top: 16px; border-bottom: 1px solid #cfd8dc;
     padding-bottom: 2px; page-break-after: avoid; }
h3 { font-size: 11.5pt; }
code { background: #eef1f3; padding: 1px 4px; border-radius: 4px; font-size: 9pt;
       font-family: 'Cascadia Code', Consolas, monospace; }
pre { background: #f4f6f8; border: 1px solid #dde3e8; border-radius: 6px; padding: 8px 10px;
      overflow-x: auto; page-break-inside: avoid; font-size: 8.6pt; }
pre code { background: none; padding: 0; }
table { border-collapse: collapse; width: 100%; margin: 6px 0; page-break-inside: avoid; }
th, td { border: 1px solid #c4ccd2; padding: 4px 7px; text-align: left; vertical-align: top; }
th { background: #eaf3ee; }
blockquote { border-left: 3px solid #2e8b57; margin: 8px 0; padding: 2px 12px; background: #f5f9f7;
             color: #333; }
a { color: #0b6; text-decoration: none; }
div.mermaid { background: transparent; border: 0; text-align: center;
              page-break-inside: avoid; padding: 4px 0; }
div.mermaid svg { max-width: 100%; height: auto; }
"""


# fenced_code emite <pre><code class="language-mermaid"> con el contenido
# escapado; de ahi se saca el codigo del diagrama para pedir su SVG.
MERMAID_HTML_RE = re.compile(
    r'<pre><code class="language-mermaid">(.*?)</code></pre>', re.DOTALL
)


def incrustar_diagramas(html_text: str) -> tuple[str, int]:
    """Sustituye cada bloque mermaid por su SVG. Devuelve (html, sin dibujar)."""
    fuentes = [htmllib.unescape(m) for m in MERMAID_HTML_RE.findall(html_text)]
    if not fuentes:
        return html_text, 0
    svgs = dibujar_diagramas(fuentes, verbose=False)
    faltan = 0

    def reemplazo(match):
        nonlocal faltan
        svg = svgs.get(htmllib.unescape(match.group(1)).strip())
        if not svg:
            faltan += 1
            return match.group(0)
        return f'<div class="mermaid">{svg}</div>'

    return MERMAID_HTML_RE.sub(reemplazo, html_text), faltan


def md_a_html(md_text: str) -> tuple[str, int]:
    cuerpo = markdown.markdown(
        md_text, extensions=["tables", "fenced_code", "sane_lists", "attr_list"]
    )
    cuerpo, faltan = incrustar_diagramas(cuerpo)
    return (
        "<!doctype html><html lang='es'><head><meta charset='utf-8'>"
        f"<style>{CSS}</style></head><body>{cuerpo}</body></html>"
    ), faltan


# Tamano por debajo del cual un PDF de clase no puede ser real (la mas corta
# ronda las 6 paginas). Sirve para no dar por buena una impresion vacia.
MIN_PDF_BYTES = 20_000


def generar_pdf(nav: str, html_path: str, pdf_path: str) -> None:
    """Imprime la pagina y **comprueba** que el PDF salio.

    Se imprime a un fichero temporal y se mueve encima del definitivo. Dos
    razones: el navegador escribe con stderr silenciado, asi que si no puede
    escribir el destino (en Windows pasa mientras el antivirus tiene el fichero
    anterior mapeado: truncarlo devuelve EINVAL) el fallo era invisible y la
    clase se quedaba con su PDF viejo como si nada; y ademas una interrupcion a
    medias ya no deja un PDF truncado en el repositorio.

    No hay JavaScript que esperar —los diagramas vienen ya dibujados—, asi que
    tampoco hace falta presupuesto de tiempo virtual.
    """
    tmp_pdf = pdf_path + ".tmp"
    if os.path.exists(tmp_pdf):
        os.remove(tmp_pdf)
    uri = "file:///" + html_path.replace("\\", "/")
    subprocess.run(
        [nav, "--headless=new", "--disable-gpu", "--no-first-run",
         "--no-default-browser-check", "--no-pdf-header-footer",
         f"--print-to-pdf={tmp_pdf}", uri],
        check=True, timeout=180,
        stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL,
    )
    publicar(tmp_pdf, pdf_path, MIN_PDF_BYTES)


# ---------- PPTX ----------
VERDE = RGBColor(0x0B, 0x3D, 0x2E)
VERDE2 = RGBColor(0x2E, 0x8B, 0x57)
GRIS = RGBColor(0x22, 0x28, 0x2E)


def limpiar_inline(s: str) -> str:
    s = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", s)  # enlaces -> texto
    s = re.sub(r"[*`_]+", "", s)                     # énfasis / code
    return s.strip()


# Los bloques cercados (```bash, ```mermaid, ```python) no son prosa: sin esto
# sus lineas acababan como vinetas del deck ("flowchart TD", "sudo nmap -sS").
FENCE_RE = re.compile(r"^[ \t]*```.*?^[ \t]*```[ \t]*$\n?", re.MULTILINE | re.DOTALL)


def sin_bloques_codigo(md_text: str) -> str:
    return FENCE_RE.sub("", md_text)


def partir_secciones(md_text: str) -> list[tuple[str, str]]:
    # separa por encabezados de nivel 2
    partes = re.split(r"\n##\s+", "\n" + sin_bloques_codigo(md_text))
    out = []
    for p in partes[1:]:
        linea, _, resto = p.partition("\n")
        out.append((linea.strip(), resto.strip()))
    return out


MERMAID_MD_RE = re.compile(r"```mermaid\n(.*?)```", re.DOTALL)


def diagramas_de(md_text: str) -> list[str]:
    """Codigo de cada diagrama de la clase, en el orden en que aparece."""
    return [b.strip() for b in MERMAID_MD_RE.findall(md_text) if b.strip()]


MAX_VINETA = 200  # una diapositiva no sostiene un parrafo entero de la clase


def recortar(texto: str, limite: int = MAX_VINETA) -> str:
    if len(texto) <= limite:
        return texto
    return texto[:limite].rsplit(" ", 1)[0].rstrip(" ,;:") + "…"


def bullets_de(cuerpo: str, maximo: int = 7) -> list[str]:
    items: list[str] = []
    for linea in cuerpo.splitlines():
        l = linea.strip()
        if not l:
            continue
        if l.startswith("|"):  # fila de tabla -> 1ª y 2ª celda
            celdas = [c.strip() for c in l.strip("|").split("|")]
            if set("".join(celdas)) <= set("-: "):
                continue
            if celdas and celdas[0] in ("#", "Tema", "Síntoma / mensaje", "Síntoma"):
                continue
            texto = " — ".join([c for c in celdas[:2] if c])
            items.append(limpiar_inline(texto))
        elif re.match(r"^(\d+\.|[-*])\s+", l):
            items.append(limpiar_inline(re.sub(r"^(\d+\.|[-*])\s+", "", l)))
        elif l.startswith(">") or l.startswith("#") or l.startswith("```"):
            continue
        elif l.startswith(":"):  # definición
            items.append(limpiar_inline(l[1:].strip()))
        else:
            if len(l) > 3:
                items.append(limpiar_inline(l))
        if len(items) >= maximo:
            break
    return [recortar(i) for i in items if i][:maximo]


def add_slide_contenido(prs, titulo: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # en blanco
    # barra de título
    izq, arr, ancho = Inches(0.5), Inches(0.35), Inches(9)
    tb = slide.shapes.add_textbox(izq, arr, ancho, Inches(0.9)).text_frame
    tb.word_wrap = True
    p = tb.paragraphs[0]
    r = p.add_run(); r.text = titulo
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = VERDE
    # cuerpo
    cuerpo = slide.shapes.add_textbox(izq, Inches(1.4), ancho, Inches(5.4)).text_frame
    cuerpo.word_wrap = True
    for i, b in enumerate(bullets):
        par = cuerpo.paragraphs[0] if i == 0 else cuerpo.add_paragraph()
        run = par.add_run(); run.text = "•  " + b
        run.font.size = Pt(15); run.font.color.rgb = GRIS
        par.space_after = Pt(6)


def add_slide_diagrama(prs, titulo: str, png: str) -> None:
    """Una diapositiva con el diagrama, centrado y a la mayor escala que quepa."""
    from PIL import Image  # solo se usa para leer el tamano real del PNG

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    caja = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7)).text_frame
    caja.word_wrap = True
    par = caja.paragraphs[0]
    run = par.add_run()
    run.text = titulo
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = VERDE

    with Image.open(png) as img:
        ancho_px, alto_px = img.size
    # Area util de la diapositiva por debajo del titulo.
    max_ancho, max_alto = Inches(9.0), Inches(5.8)
    escala = min(max_ancho / ancho_px, max_alto / alto_px)
    ancho, alto = int(ancho_px * escala), int(alto_px * escala)
    izq = int((prs.slide_width - ancho) / 2)
    arr = Inches(1.2) + int((max_alto - alto) / 2)
    slide.shapes.add_picture(png, izq, arr, width=ancho, height=alto)


def construir_pptx(md_text: str, titulo: str, subtitulo: str, out_path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    # portada
    s = prs.slides.add_slide(prs.slide_layouts[6])
    caja = s.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(8.8), Inches(2.5)).text_frame
    caja.word_wrap = True
    p = caja.paragraphs[0]; r = p.add_run(); r.text = titulo
    r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = VERDE
    p2 = caja.add_paragraph(); r2 = p2.add_run(); r2.text = subtitulo
    r2.font.size = Pt(16); r2.font.color.rgb = VERDE2
    p3 = caja.add_paragraph(); r3 = p3.add_run()
    r3.text = "Programa de Ciberseguridad Moderna"
    r3.font.size = Pt(12); r3.font.color.rgb = GRIS

    saltar = {"⬅️ Clase anterior", "➡️ Siguiente clase", "📥 Material descargable"}
    for titulo_sec, cuerpo in partir_secciones(md_text):
        if any(k in titulo_sec for k in saltar):
            continue
        bullets = bullets_de(cuerpo)
        if not bullets:
            continue
        add_slide_contenido(prs, titulo_sec, bullets)

    # Los diagramas, al final: el deck se proyecta y el grafico es lo que sostiene
    # la explicacion. Antes no viajaba ninguno, porque partir_secciones descarta
    # los bloques de codigo y el diagrama es uno de ellos.
    fuentes = diagramas_de(md_text)
    if fuentes:
        pngs = rasterizar_diagramas(fuentes, verbose=False)
        vistos: set[str] = set()
        for numero, fuente in enumerate(fuentes, start=1):
            png = pngs.get(fuente)
            if not png or fuente in vistos:
                continue
            vistos.add(fuente)
            etiqueta = "Diagrama de la clase" if len(fuentes) == 1 else f"Diagrama {numero}"
            add_slide_diagrama(prs, etiqueta, str(png))
    # Igual que con el PDF: se guarda aparte y se mueve encima, para no
    # quedarse sin presentacion si el destino no se puede truncar.
    tmp_path = out_path + ".tmp"
    prs.save(tmp_path)
    publicar(tmp_path, out_path)


# ---------- README: sección de material ----------
def anadir_seccion_descargas(readme: str, pdf_name: str, pptx_name: str) -> None:
    txt = open(readme, encoding="utf-8").read()
    if "## 📥 Material descargable" in txt:
        return
    bloque = (
        "## 📥 Material descargable\n\n"
        f"- 📄 [Guía en PDF](./{pdf_name}) — versión imprimible de esta clase.\n"
        f"- 🎞️ [Presentación (PPTX)](./{pptx_name}) — deck para proyectar en clase.\n\n"
    )
    marcador = "## ➡️ Siguiente clase"
    if marcador in txt:
        txt = txt.replace(marcador, bloque + marcador, 1)
    else:
        txt = txt.rstrip() + "\n\n" + bloque
    open(readme, "w", encoding="utf-8", newline="\n").write(txt)


def indices_de_partes() -> list[int]:
    """Numeros de parte presentes en classes/, en orden."""
    nums = []
    for d in glob.glob(os.path.join(CLASSES, "parte-*")):
        m = re.match(r"parte-(\d+)", os.path.basename(d))
        if m and os.path.isdir(d):
            nums.append(int(m.group(1)))
    return sorted(nums)


def generar_parte(idx: int, solo_una: bool = False) -> int:
    partes = sorted(glob.glob(os.path.join(CLASSES, f"parte-{idx}-*")))
    if not partes:
        raise SystemExit(f"No existe la parte {idx}")
    pdir = partes[0]
    clases = sorted(d for d in glob.glob(os.path.join(pdir, "*")) if os.path.isdir(d))
    if solo_una:
        clases = clases[:1]

    nav = navegador()
    tmp = tempfile.mkdtemp(prefix="matcurso_")
    hechos = 0
    for cdir in clases:
        readme = os.path.join(cdir, "README.md")
        if not os.path.isfile(readme):
            continue
        md_text = open(readme, encoding="utf-8").read()
        slug = os.path.basename(cdir)
        m = re.search(r"^#\s+(.+)$", md_text, re.MULTILINE)
        titulo = re.sub(r"[#*`]", "", m.group(1)).strip() if m else slug
        num = slug[:3]
        pdf_name = f"clase-{num}-guia.pdf"
        pptx_name = f"clase-{num}-presentacion.pptx"

        # PDF
        html_path = os.path.join(tmp, f"{slug}.html")
        html_text, sin_dibujar = md_a_html(md_text)
        if sin_dibujar:
            print(f"  AVISO {slug}: {sin_dibujar} diagrama(s) sin dibujar", flush=True)
        open(html_path, "w", encoding="utf-8").write(html_text)
        generar_pdf(nav, html_path, os.path.join(cdir, pdf_name))
        # PPTX
        construir_pptx(md_text, titulo, os.path.basename(pdir).replace("-", " "),
                       os.path.join(cdir, pptx_name))
        # README
        anadir_seccion_descargas(readme, pdf_name, pptx_name)
        hechos += 1
        print(f"  [OK] {slug}  -> {pdf_name} + {pptx_name}", flush=True)

    print(f"Parte {idx}: material generado para {hechos} clase(s).", flush=True)
    return hechos


def main() -> int:
    if "--todas" in sys.argv:
        total = sum(generar_parte(i) for i in indices_de_partes())
        print(f"TOTAL: material generado para {total} clase(s).")
        return 0
    if len(sys.argv) < 2:
        raise SystemExit(
            "Uso: python scripts/generar_material.py <indice_parte> [--solo-una]\n"
            "     python scripts/generar_material.py --todas"
        )
    generar_parte(int(sys.argv[1]), "--solo-una" in sys.argv)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
